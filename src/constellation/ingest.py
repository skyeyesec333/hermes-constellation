"""Source-preserving, deterministic local ingest."""

from __future__ import annotations

import importlib.metadata
import importlib.util
import io
import json
import mimetypes
import platform
import zipfile
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import Path, PurePosixPath
from typing import Any, cast
from urllib.parse import parse_qsl, urlsplit

import yaml

from .card_ingest import extract_business_card_fields
from .frontmatter import parse_frontmatter, render_frontmatter
from .models import CandidatePatch, Sensitivity, SourceItem
from .storage import (
    atomic_write_bytes,
    atomic_write_text,
    safe_relative_path,
    sha256_bytes,
    sha256_file,
)
from .vault import CONFIG_RELATIVE, is_initialized

_ULID_ALPHABET = "0123456789ABCDEFGHJKMNPQRSTVWXYZ"
MAX_SOURCE_BYTES = 50 * 1024 * 1024
MAX_EXTRACTED_CHARS = 5_000_000
MAX_PDF_PAGES = 500
MAX_IMAGE_PIXELS = 40_000_000
MAX_OFFICE_UNCOMPRESSED_BYTES = 200 * 1024 * 1024
MAX_OFFICE_ARCHIVE_ENTRIES = 10_000
MAX_OFFICE_COMPRESSION_RATIO = 1_000
OCR_MIN_CONFIDENCE = 0.45
OCR_MIN_CHARACTERS = 3

OOXML_MAIN_CONTENT_TYPES = {
    ".docx": (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"
    ),
    ".pptx": (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
    ),
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml",
}


class IngestError(RuntimeError):
    pass


class CapabilityError(IngestError):
    pass


def _source_registration_mode(vault: Path) -> str:
    config_path = safe_relative_path(vault, CONFIG_RELATIVE)
    try:
        config = yaml.safe_load(config_path.read_text(encoding="utf-8"))
    except (OSError, yaml.YAMLError) as exc:
        raise IngestError("vault configuration is unreadable") from exc
    mode = config.get("source_registration", "review") if isinstance(config, dict) else None
    if mode not in {"review", "automatic"}:
        raise IngestError("source_registration must be review or automatic")
    return mode


def _validated_source_url(source_url: str | None) -> str | None:
    if source_url is None:
        return None
    parsed = urlsplit(source_url)
    if parsed.scheme not in {"http", "https"} or not parsed.netloc:
        raise IngestError("source URL must be an absolute http or https URL")
    if parsed.username or parsed.password:
        raise IngestError("source URL must not contain credentials")
    sensitive_query_keys = {
        "access_token",
        "api_key",
        "apikey",
        "authorization",
        "password",
        "secret",
        "signature",
        "sig",
        "token",
    }
    if any(
        key.casefold() in sensitive_query_keys or key.casefold().startswith("x-amz-")
        for key, _ in parse_qsl(parsed.query, keep_blank_values=True)
    ):
        raise IngestError("source URL must not contain credentials")
    return source_url


@dataclass(frozen=True)
class ExtractedSource:
    data: bytes
    text: str
    media_type: str
    extraction: dict[str, Any]


def _id_from_hash(digest: str) -> str:
    value = int(digest[:32], 16)
    characters = ["0"] * 26
    for index in range(25, -1, -1):
        characters[index] = _ULID_ALPHABET[value & 31]
        value >>= 5
    return "".join(characters)


def _line_count(text: str) -> int:
    return len(text.splitlines()) or 1


def _text_extraction(data: bytes, text: str, media_type: str) -> ExtractedSource:
    if not text.strip():
        raise IngestError("text source contains no extractable content")
    lines = _line_count(text)
    return ExtractedSource(
        data=data,
        text=text,
        media_type=media_type,
        extraction={
            "schema_version": "0.1",
            "status": "complete",
            "engine": {
                "name": "python-utf8",
                "version": platform.python_version(),
                "options": {"encoding": "utf-8", "errors": "strict"},
            },
            "source_sha256": sha256_bytes(data),
            "extracted_text_sha256": sha256_bytes(text.encode("utf-8")),
            "characters": len(text),
            "expected_units": 1,
            "extracted_units": 1,
            "blank_units": 0,
            "failed_units": 0,
            "truncated_units": 0,
            "warnings": [],
            "units": [
                {
                    "kind": "text",
                    "index": 1,
                    "status": "extracted",
                    "anchor": f"L000001-L{lines:06d}",
                    "line_start": 1,
                    "line_end": lines,
                    "characters": len(text),
                    "text_sha256": sha256_bytes(text.encode("utf-8")),
                }
            ],
        },
    )


def _detect_media_type(data: bytes) -> str:
    try:
        import magic  # type: ignore[import-not-found]
    except ImportError as exc:
        raise CapabilityError("MIME detection requires the python-magic capability") from exc

    try:
        detected = str(magic.from_buffer(data, mime=True)).strip().lower()
    except Exception as exc:
        raise IngestError("MIME detection failed") from exc
    if not detected:
        raise IngestError("MIME detection returned no media type")
    return detected


def _validate_ooxml_archive(data: bytes, suffix: str) -> None:
    expected_content_type = OOXML_MAIN_CONTENT_TYPES[suffix]
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            members = archive.infolist()
            if len(members) > MAX_OFFICE_ARCHIVE_ENTRIES:
                raise IngestError("OOXML archive exceeds the configured entry limit")
            expanded_bytes = 0
            for member in members:
                normalized = PurePosixPath(member.filename.replace("\\", "/"))
                if normalized.is_absolute() or ".." in normalized.parts:
                    raise IngestError("OOXML archive contains an unsafe member path")
                if member.flag_bits & 0x1:
                    raise IngestError("encrypted OOXML archives are not supported")
                expanded_bytes += member.file_size
                if expanded_bytes > MAX_OFFICE_UNCOMPRESSED_BYTES:
                    raise IngestError("OOXML archive exceeds the configured expanded size")
                ratio = member.file_size / max(member.compress_size, 1)
                if member.file_size > 1_000_000 and ratio > MAX_OFFICE_COMPRESSION_RATIO:
                    raise IngestError("OOXML archive contains a suspicious compression ratio")
            try:
                content_types = archive.read("[Content_Types].xml")
            except KeyError as exc:
                raise IngestError("OOXML archive is missing its content-type manifest") from exc
    except IngestError:
        raise
    except zipfile.BadZipFile as exc:
        raise IngestError("OOXML source is not a valid ZIP archive") from exc
    if expected_content_type.encode("ascii") not in content_types:
        raise IngestError("file extension and OOXML content type do not match")


def _docx_extraction(data: bytes) -> ExtractedSource:
    if importlib.util.find_spec("docx") is None:
        raise CapabilityError("DOCX extraction requires the optional python-docx capability")
    import docx  # type: ignore[import-not-found]

    try:
        document = docx.Document(io.BytesIO(data))
        sections: list[str] = []
        units: list[dict[str, Any]] = []
        paragraph_index = 0
        table_index = 0
        for block in document.iter_inner_content():
            if isinstance(block, docx.text.paragraph.Paragraph):
                value = block.text.strip()
                if not value:
                    continue
                paragraph_index += 1
                anchor = f"PARA{paragraph_index:04d}"
                sections.append(f"[{anchor}] {value}")
                units.append(
                    {
                        "kind": "paragraph",
                        "index": paragraph_index,
                        "status": "extracted",
                        "anchor": anchor,
                        "style": block.style.name if block.style else None,
                        "characters": len(value),
                        "text_sha256": sha256_bytes(value.encode("utf-8")),
                    }
                )
                continue
            table_index += 1
            for row_index, row in enumerate(block.rows, start=1):
                for column_index, cell in enumerate(row.cells, start=1):
                    value = cell.text.strip()
                    if not value:
                        continue
                    anchor = f"TABLE{table_index:04d}:R{row_index:04d}:C{column_index:04d}"
                    sections.append(f"[{anchor}] {value}")
                    units.append(
                        {
                            "kind": "table-cell",
                            "index": len(units) + 1,
                            "status": "extracted",
                            "anchor": anchor,
                            "table": table_index,
                            "row": row_index,
                            "column": column_index,
                            "characters": len(value),
                            "text_sha256": sha256_bytes(value.encode("utf-8")),
                        }
                    )
        if not sections:
            raise IngestError("DOCX contains no extractable text")
        text = "\n".join(sections) + "\n"
        if len(text) > MAX_EXTRACTED_CHARS:
            raise IngestError("extracted text exceeds the configured size limit")
        return ExtractedSource(
            data=data,
            text=text,
            media_type=(
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ),
            extraction={
                "schema_version": "0.1",
                "status": "complete",
                "engine": {
                    "name": "python-docx",
                    "version": str(getattr(docx, "__version__", "unknown")),
                    "options": {"method": "iter_inner_content"},
                },
                "source_sha256": sha256_bytes(data),
                "extracted_text_sha256": sha256_bytes(text.encode("utf-8")),
                "characters": len(text),
                "expected_units": len(units),
                "extracted_units": len(units),
                "blank_units": 0,
                "failed_units": 0,
                "truncated_units": 0,
                "warnings": [],
                "units": units,
            },
        )
    except (IngestError, CapabilityError):
        raise
    except Exception as exc:
        raise IngestError("DOCX extraction failed") from exc


def _pptx_extraction(data: bytes) -> ExtractedSource:
    if importlib.util.find_spec("pptx") is None:
        raise CapabilityError("PPTX extraction requires the optional python-pptx capability")
    import pptx  # type: ignore[import-not-found]

    try:
        presentation = pptx.Presentation(io.BytesIO(data))
        sections: list[str] = []
        units: list[dict[str, Any]] = []
        blank_units = 0
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_lines: list[str] = []
            text_index = 0
            table_index = 0
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    value = shape.text.strip()
                    if value:
                        text_index += 1
                        anchor = f"SLIDE{slide_index:04d}:TEXT{text_index:04d}"
                        slide_lines.append(f"[{anchor}] {value}")
                if getattr(shape, "has_table", False):
                    table_index += 1
                    for row_index, row in enumerate(shape.table.rows, start=1):
                        for column_index, cell in enumerate(row.cells, start=1):
                            value = cell.text.strip()
                            if not value:
                                continue
                            anchor = (
                                f"SLIDE{slide_index:04d}:TABLE{table_index:04d}:"
                                f"R{row_index:04d}:C{column_index:04d}"
                            )
                            slide_lines.append(f"[{anchor}] {value}")
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_lines.append(f"[SLIDE{slide_index:04d}:NOTES] {notes_text}")
            slide_text = "\n".join(slide_lines)
            if slide_text:
                sections.append(slide_text)
                status = "extracted"
            else:
                blank_units += 1
                sections.append(f"[SLIDE{slide_index:04d}]")
                status = "blank"
            units.append(
                {
                    "kind": "slide",
                    "index": slide_index,
                    "status": status,
                    "anchor": f"SLIDE{slide_index:04d}",
                    "notes": bool(notes_text),
                    "tables": table_index,
                    "characters": len(slide_text),
                    "text_sha256": sha256_bytes(slide_text.encode("utf-8")),
                }
            )
        if not units or blank_units == len(units):
            raise IngestError("PPTX contains no extractable text")
        text = "\n\n".join(sections) + "\n"
        if len(text) > MAX_EXTRACTED_CHARS:
            raise IngestError("extracted text exceeds the configured size limit")
        status = "complete-with-gaps" if blank_units else "complete"
        warnings = [f"{blank_units} slide(s) contain no extractable text"] if blank_units else []
        return ExtractedSource(
            data=data,
            text=text,
            media_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
            extraction={
                "schema_version": "0.1",
                "status": status,
                "engine": {
                    "name": "python-pptx",
                    "version": str(getattr(pptx, "__version__", "unknown")),
                    "options": {"include": ["text", "tables", "speaker-notes"]},
                },
                "source_sha256": sha256_bytes(data),
                "extracted_text_sha256": sha256_bytes(text.encode("utf-8")),
                "characters": len(text),
                "expected_units": len(units),
                "extracted_units": len(units) - blank_units,
                "blank_units": blank_units,
                "failed_units": 0,
                "truncated_units": 0,
                "warnings": warnings,
                "units": units,
            },
        )
    except (IngestError, CapabilityError):
        raise
    except Exception as exc:
        raise IngestError("PPTX extraction failed") from exc


def _xlsx_extraction(data: bytes) -> ExtractedSource:
    if importlib.util.find_spec("openpyxl") is None:
        raise CapabilityError("XLSX extraction requires the optional openpyxl capability")
    import openpyxl  # type: ignore[import-not-found]

    workbook = None
    try:
        workbook = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=False)
        sections: list[str] = []
        units: list[dict[str, Any]] = []
        for sheet_index, worksheet in enumerate(workbook.worksheets, start=1):
            sheet_anchor = f"SHEET{sheet_index:04d}"
            sections.append(f"[{sheet_anchor}] {worksheet.title}")
            for row in worksheet.iter_rows():
                for cell in row:
                    if cell.value is None:
                        continue
                    serializer = getattr(cell.value, "isoformat", None)
                    value = str(serializer()) if callable(serializer) else str(cell.value)
                    anchor = f"{sheet_anchor}:{cell.coordinate}"
                    sections.append(f"[{anchor}] {value}")
                    units.append(
                        {
                            "kind": "cell",
                            "index": len(units) + 1,
                            "status": "extracted",
                            "anchor": anchor,
                            "sheet": worksheet.title,
                            "coordinate": cell.coordinate,
                            "formula": isinstance(cell.value, str) and cell.value.startswith("="),
                            "characters": len(value),
                            "text_sha256": sha256_bytes(value.encode("utf-8")),
                        }
                    )
        if not units:
            raise IngestError("XLSX contains no extractable cells")
        text = "\n".join(sections) + "\n"
        if len(text) > MAX_EXTRACTED_CHARS:
            raise IngestError("extracted text exceeds the configured size limit")
        return ExtractedSource(
            data=data,
            text=text,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            extraction={
                "schema_version": "0.1",
                "status": "complete",
                "engine": {
                    "name": "openpyxl",
                    "version": str(getattr(openpyxl, "__version__", "unknown")),
                    "options": {"read_only": True, "data_only": False},
                },
                "source_sha256": sha256_bytes(data),
                "extracted_text_sha256": sha256_bytes(text.encode("utf-8")),
                "characters": len(text),
                "expected_units": len(units),
                "extracted_units": len(units),
                "blank_units": 0,
                "failed_units": 0,
                "truncated_units": 0,
                "warnings": [],
                "units": units,
            },
        )
    except (IngestError, CapabilityError):
        raise
    except Exception as exc:
        raise IngestError("XLSX extraction failed") from exc
    finally:
        if workbook is not None:
            workbook.close()


def _rapidocr_engine() -> Any:
    if importlib.util.find_spec("rapidocr_onnxruntime") is None:
        raise CapabilityError("OCR requires the optional RapidOCR capability")
    from rapidocr_onnxruntime import RapidOCR  # type: ignore[import-not-found]

    return RapidOCR()


def _rapidocr_regions(data: bytes, *, engine: Any | None = None) -> list[dict[str, Any]]:
    active_engine = engine or _rapidocr_engine()
    raw_result, _ = active_engine(data)
    regions: list[dict[str, Any]] = []
    for item in raw_result or []:
        box, raw_text, raw_confidence = item
        value = str(raw_text).strip()
        confidence = float(raw_confidence)
        if not value or confidence < OCR_MIN_CONFIDENCE:
            continue
        regions.append(
            {
                "text": value,
                "confidence": round(confidence, 6),
                "bounding_box": [
                    [round(float(coordinate), 3) for coordinate in point] for point in box
                ],
            }
        )
    if sum(len(region["text"]) for region in regions) < OCR_MIN_CHARACTERS:
        raise CapabilityError("OCR produced no reliable text; use the configured vision fallback")
    return regions


def _image_extraction(data: bytes, suffix: str) -> ExtractedSource:
    if importlib.util.find_spec("PIL") is None:
        raise CapabilityError("image extraction requires the optional Pillow capability")
    from PIL import Image  # type: ignore[import-not-found]

    try:
        with Image.open(io.BytesIO(data)) as image:
            width, height = image.size
            detected_format = str(image.format or "unknown").upper()
            expected_format = {
                ".png": "PNG",
                ".jpg": "JPEG",
                ".jpeg": "JPEG",
                ".webp": "WEBP",
                ".tif": "TIFF",
                ".tiff": "TIFF",
                ".bmp": "BMP",
            }[suffix]
            if detected_format != expected_format:
                raise IngestError("file extension and image signature do not match")
            if width * height > MAX_IMAGE_PIXELS:
                raise IngestError("image exceeds the configured pixel limit")
            image.verify()
        regions = _rapidocr_regions(data)
        units: list[dict[str, Any]] = []
        sections: list[str] = []
        for index, region in enumerate(regions, start=1):
            value = str(region["text"])
            anchor = f"OCR:R{index:04d}"
            sections.append(f"[{anchor}] {value}")
            units.append(
                {
                    "kind": "ocr-region",
                    "index": index,
                    "status": "extracted",
                    "anchor": anchor,
                    "confidence": region["confidence"],
                    "bounding_box": region["bounding_box"],
                    "characters": len(value),
                    "text_sha256": sha256_bytes(value.encode("utf-8")),
                }
            )
        text = "\n".join(sections) + "\n"
        average_confidence = round(
            sum(float(unit["confidence"]) for unit in units) / len(units), 6
        )
        warnings = []
        status = "complete"
        if average_confidence < 0.65:
            status = "complete-low-confidence"
            warnings.append("average OCR confidence is below 0.65; verify with vision")
        media_type = mimetypes.types_map.get(suffix, f"image/{detected_format.lower()}")
        return ExtractedSource(
            data=data,
            text=text,
            media_type=media_type,
            extraction={
                "schema_version": "0.1",
                "status": status,
                "engine": {
                    "name": "RapidOCR",
                    "version": importlib.metadata.version("rapidocr-onnxruntime"),
                    "options": {
                        "minimum_confidence": OCR_MIN_CONFIDENCE,
                        "minimum_characters": OCR_MIN_CHARACTERS,
                    },
                },
                "source_sha256": sha256_bytes(data),
                "extracted_text_sha256": sha256_bytes(text.encode("utf-8")),
                "characters": len(text),
                "image": {"format": detected_format, "width": width, "height": height},
                "average_confidence": average_confidence,
                "expected_units": len(units),
                "extracted_units": len(units),
                "blank_units": 0,
                "failed_units": 0,
                "truncated_units": 0,
                "warnings": warnings,
                "units": units,
            },
        )
    except (IngestError, CapabilityError):
        raise
    except Exception as exc:
        raise IngestError("image OCR extraction failed") from exc


def _pdf_extraction(data: bytes) -> ExtractedSource:
    if importlib.util.find_spec("fitz") is None:
        raise CapabilityError("PDF extraction requires the optional PyMuPDF capability")
    import fitz  # type: ignore[import-not-found]

    document = None
    try:
        document = fitz.open(stream=data, filetype="pdf")
        if document.page_count > MAX_PDF_PAGES:
            raise IngestError("PDF exceeds the configured page limit")
        if document.page_count < 1:
            raise IngestError("PDF contains no pages")
        sections: list[str] = []
        units: list[dict[str, Any]] = []
        warnings: list[str] = []
        extracted_units = 0
        blank_units = 0
        ocr_attempted = 0
        ocr_extracted = 0
        ocr_confidences: list[float] = []
        ocr_engine: Any | None = None
        for page_index in range(1, document.page_count + 1):
            page = document.load_page(page_index - 1)
            native_text = cast(str, page.get_text("text"))
            marker = f"[P{page_index:04d}]"
            if native_text.strip():
                page_text = native_text.rstrip()
                sections.append(f"{marker}\n{page_text}")
                extracted_units += 1
                lines = _line_count(native_text)
                units.append(
                    {
                        "kind": "page",
                        "index": page_index,
                        "status": "extracted",
                        "method": "native-text",
                        "anchor": f"P{page_index:04d}:L0001-L{lines:04d}",
                        "line_start": 1,
                        "line_end": lines,
                        "regions": [],
                        "characters": len(native_text),
                        "text_sha256": sha256_bytes(native_text.encode("utf-8")),
                    }
                )
                continue

            ocr_attempted += 1
            try:
                if ocr_engine is None:
                    ocr_engine = _rapidocr_engine()
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                detected = _rapidocr_regions(pixmap.tobytes("png"), engine=ocr_engine)
                regions: list[dict[str, Any]] = []
                page_lines: list[str] = []
                for region_index, region in enumerate(detected, start=1):
                    anchor = f"P{page_index:04d}:OCR:R{region_index:04d}"
                    value = str(region["text"])
                    page_lines.append(f"[{anchor}] {value}")
                    regions.append({"anchor": anchor, **region})
                    ocr_confidences.append(float(region["confidence"]))
                page_text = "\n".join(page_lines)
                average_confidence = round(
                    sum(float(region["confidence"]) for region in regions) / len(regions), 6
                )
                sections.append(f"{marker}\n{page_text}")
                extracted_units += 1
                ocr_extracted += 1
                units.append(
                    {
                        "kind": "page",
                        "index": page_index,
                        "status": "ocr-extracted",
                        "method": "rapidocr",
                        "anchor": f"P{page_index:04d}:OCR",
                        "line_start": None,
                        "line_end": None,
                        "regions": regions,
                        "average_confidence": average_confidence,
                        "characters": sum(len(str(region["text"])) for region in regions),
                        "text_sha256": sha256_bytes(page_text.encode("utf-8")),
                    }
                )
            except CapabilityError:
                blank_units += 1
                sections.append(marker)
                units.append(
                    {
                        "kind": "page",
                        "index": page_index,
                        "status": "blank-needs-vision",
                        "method": "none",
                        "anchor": f"P{page_index:04d}",
                        "line_start": None,
                        "line_end": None,
                        "regions": [],
                        "characters": 0,
                        "text_sha256": sha256_bytes(b""),
                    }
                )
        if extracted_units == 0:
            raise CapabilityError(
                "PDF contains no reliable native or OCR text; requires OCR or vision fallback"
            )
        text = "\n\n".join(sections) + "\n"
        if len(text) > MAX_EXTRACTED_CHARS:
            raise IngestError("extracted text exceeds the configured size limit")
        status = "complete"
        if blank_units:
            status = "complete-with-gaps"
            warnings.append(f"{blank_units} page(s) need vision verification")
        average_ocr_confidence = (
            round(sum(ocr_confidences) / len(ocr_confidences), 6)
            if ocr_confidences
            else None
        )
        if average_ocr_confidence is not None and average_ocr_confidence < 0.65:
            if not blank_units:
                status = "complete-low-confidence"
            warnings.append("average PDF OCR confidence is below 0.65; verify with vision")
        version = str(getattr(fitz, "VersionBind", getattr(fitz, "__version__", "unknown")))
        engine_name = "PyMuPDF+RapidOCR" if ocr_extracted else "PyMuPDF"
        return ExtractedSource(
            data=data,
            text=text,
            media_type="application/pdf",
            extraction={
                "schema_version": "0.1",
                "status": status,
                "engine": {
                    "name": engine_name,
                    "version": version,
                    "ocr_version": (
                        importlib.metadata.version("rapidocr-onnxruntime")
                        if ocr_extracted
                        else None
                    ),
                    "options": {
                        "native_method": "page.get_text",
                        "native_mode": "text",
                        "ocr_render_scale": 2,
                        "ocr_minimum_confidence": OCR_MIN_CONFIDENCE,
                    },
                },
                "source_sha256": sha256_bytes(data),
                "extracted_text_sha256": sha256_bytes(text.encode("utf-8")),
                "characters": len(text),
                "expected_units": document.page_count,
                "extracted_units": extracted_units,
                "blank_units": blank_units,
                "failed_units": 0,
                "truncated_units": 0,
                "ocr": {
                    "attempted_pages": ocr_attempted,
                    "extracted_pages": ocr_extracted,
                    "average_confidence": average_ocr_confidence,
                },
                "warnings": warnings,
                "units": units,
            },
        )
    except (IngestError, CapabilityError):
        raise
    except Exception as exc:
        raise IngestError("PDF extraction failed") from exc
    finally:
        if document is not None:
            document.close()


def _read_source(path: Path) -> ExtractedSource:
    suffix = path.suffix.lower()
    if path.stat().st_size > MAX_SOURCE_BYTES:
        raise IngestError("source exceeds the configured size limit")
    data = path.read_bytes()
    if len(data) > MAX_SOURCE_BYTES:
        raise IngestError("source exceeds the configured size limit")
    detected_media_type = _detect_media_type(data)
    if suffix in OOXML_MAIN_CONTENT_TYPES:
        _validate_ooxml_archive(data, suffix)
    if suffix in {".txt", ".md", ".markdown"}:
        if not detected_media_type.startswith("text/"):
            raise IngestError("file extension and detected text media type do not match")
        if b"\x00" in data:
            raise IngestError("text source contains binary NUL bytes")
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError as exc:
            raise IngestError("text sources must be UTF-8") from exc
        if len(text) > MAX_EXTRACTED_CHARS:
            raise IngestError("extracted text exceeds the configured size limit")
        media_type = "text/markdown" if suffix in {".md", ".markdown"} else "text/plain"
        extracted = _text_extraction(data, text, media_type)
    elif suffix == ".pdf":
        if detected_media_type != "application/pdf" or not data.startswith(b"%PDF-"):
            raise IngestError("file extension and PDF signature do not match")
        extracted = _pdf_extraction(data)
    elif suffix == ".docx":
        extracted = _docx_extraction(data)
    elif suffix == ".pptx":
        extracted = _pptx_extraction(data)
    elif suffix == ".xlsx":
        extracted = _xlsx_extraction(data)
    elif suffix in {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff", ".bmp"}:
        if not detected_media_type.startswith("image/"):
            raise IngestError("file extension and detected image media type do not match")
        extracted = _image_extraction(data, suffix)
    else:
        guessed, _ = mimetypes.guess_type(path.name)
        raise IngestError(f"unsupported source type: {guessed or suffix or 'unknown'}")
    extracted.extraction["detected_media_type"] = detected_media_type
    return extracted


def _stage_source_extraction_upgrade(
    vault: Path,
    manifest: dict[str, Any],
    manifest_relative: Path,
    extraction: dict[str, Any],
    text: str,
    instant: datetime,
) -> bool:
    """Stage a hash-checked source-note update for an already-promoted source."""
    source_item_relative = str(manifest["source_item_path"])
    source_item_path = safe_relative_path(vault, source_item_relative)
    if not source_item_path.is_file() or source_item_path.is_symlink():
        return False
    current_text = source_item_path.read_text(encoding="utf-8")
    metadata, _ = parse_frontmatter(current_text)
    metadata["extraction_manifest_path"] = manifest_relative.as_posix()
    metadata["extraction_status"] = str(extraction["status"])
    metadata["updated_at"] = instant.isoformat()
    source_item = SourceItem.model_validate(metadata, strict=False)
    updated_note = render_frontmatter(
        source_item.model_dump(mode="json", exclude_none=True),
        f"# {source_item.title}\n\n{text}",
    )
    if updated_note == current_text:
        return False
    candidate = CandidatePatch(
        id=source_item.id,
        type="candidate-patch",
        title=f"Upgrade extraction record: {source_item.title}",
        status="pending-review",
        sensitivity=source_item.sensitivity,
        created_at=instant,
        updated_at=instant,
        target_path=source_item_relative,
        content=updated_note,
        expected_base_hash=sha256_file(source_item_path),
    )
    candidate_relative = Path(".constellation/candidates") / f"{candidate.id}.json"
    candidate_path = vault / candidate_relative
    serialized = candidate.model_dump_json(indent=2) + "\n"
    if candidate_path.exists():
        try:
            existing = CandidatePatch.model_validate_json(candidate_path.read_text(encoding="utf-8"))
        except Exception as exc:
            raise IngestError("existing source candidate is invalid") from exc
        if (
            existing.target_path != source_item_relative
            or existing.expected_base_hash != candidate.expected_base_hash
            or not existing.title.startswith("Upgrade extraction record:")
        ):
            raise IngestError("a different source upgrade candidate already exists")
    else:
        atomic_write_text(vault, candidate_relative, serialized)
    manifest["candidate_id"] = candidate.id
    manifest["candidate_path"] = candidate_relative.as_posix()
    return True


def ingest_file(
    root: Path | str,
    source: Path | str,
    *,
    now: datetime | None = None,
    sensitivity: Sensitivity = Sensitivity.INTERNAL,
    source_url: str | None = None,
    kind: str = "generic",
    phone_region: str | None = None,
) -> dict[str, str]:
    """Ingest one local file; an optional URL is provenance, never fetched evidence."""
    if kind not in {"generic", "business-card"}:
        raise IngestError("ingest kind is not supported")
    vault = Path(root).absolute()
    source_url = _validated_source_url(source_url)
    if not is_initialized(vault):
        raise IngestError("vault is not initialized")
    registration_mode = _source_registration_mode(vault)
    source_value = Path(source)
    if source_value.is_absolute():
        try:
            relative = source_value.relative_to(vault)
        except ValueError as exc:
            raise IngestError("source must be inside the vault root") from exc
    else:
        relative = source_value
    try:
        source_path = safe_relative_path(vault, relative)
    except Exception as exc:
        raise IngestError("source path is unsafe") from exc
    if source_path.is_symlink() or not source_path.is_file():
        raise IngestError("source must be a regular non-symlink file")
    extracted = _read_source(source_path)
    data, text, media_type = extracted.data, extracted.text, extracted.media_type
    digest = sha256_bytes(data)
    if extracted.extraction["source_sha256"] != digest:
        raise IngestError("extraction source hash does not match preserved bytes")
    source_id = _id_from_hash(digest)
    business_card = (
        extract_business_card_fields(
            source_id=source_id,
            text=text,
            units=list(extracted.extraction.get("units", [])),
            phone_region=phone_region,
        )
        if kind == "business-card"
        else None
    )
    manifest_relative = Path(".constellation/manifests") / f"{digest}.json"
    manifest_path = vault / manifest_relative
    instant = now or datetime.now(UTC)
    if instant.tzinfo is None or instant.utcoffset() is None:
        raise IngestError("ingest timestamp must include a timezone")
    extraction = dict(extracted.extraction)
    extraction["extracted_at"] = instant.isoformat()
    if manifest_path.exists():
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        if manifest.get("source_hash") != digest:
            raise IngestError("existing manifest does not match source bytes")
        candidate_path = str(manifest["candidate_path"])
        card_updated = business_card is not None and manifest.get("business_card") != business_card
        if card_updated:
            manifest["business_card"] = business_card
        upgraded = "extraction" not in manifest
        if upgraded:
            atomic_write_text(vault, str(manifest["text_path"]), text)
            manifest.update(
                {
                    "source_size_bytes": len(data),
                    "media_type": media_type,
                    "extraction": extraction,
                }
            )
        source_patch_staged = _stage_source_extraction_upgrade(
            vault,
            manifest,
            manifest_relative,
            manifest.get("extraction", extraction),
            text,
            instant,
        )
        candidate_path = str(manifest["candidate_path"])
        if upgraded or source_patch_staged or card_updated:
            atomic_write_text(
                vault,
                manifest_relative,
                json.dumps(manifest, indent=2, sort_keys=True) + "\n",
            )
        return {
            "schema_version": "0.1",
            "status": "already_ingested",
            "source_id": str(manifest["source_id"]),
            "manifest_path": manifest_relative.as_posix(),
            "preserved_path": str(manifest["preserved_path"]),
            "text_path": str(manifest["text_path"]),
            "source_item_path": str(manifest["source_item_path"]),
            "candidate_id": str(manifest.get("candidate_id") or Path(candidate_path).stem),
            "candidate_path": candidate_path,
            "extraction_status": str(manifest.get("extraction", extraction)["status"]),
            "manifest_upgraded": str(upgraded).lower(),
            "source_patch_staged": str(source_patch_staged).lower(),
            **(
                {"business_card_fields": str(len(business_card["fields"]))}
                if business_card is not None
                else {}
            ),
        }
    preserved_relative = Path("Library/Files") / str(instant.year) / source_id / source_path.name
    text_relative = Path("Library/Text") / f"{source_id}.txt"
    source_item_relative = Path("source-items") / f"{source_id}.md"
    atomic_write_bytes(vault, preserved_relative, data)
    atomic_write_text(vault, text_relative, text)
    item = SourceItem(
        id=source_id,
        type="source-item",
        title=source_path.stem,
        status="active",
        sensitivity=sensitivity,
        created_at=instant,
        updated_at=instant,
        source_hash=digest,
        original_path=preserved_relative.as_posix(),
        extracted_text_path=text_relative.as_posix(),
        extraction_manifest_path=manifest_relative.as_posix(),
        extraction_status=extraction["status"],
        media_type=media_type,
        source_url=source_url,
    )
    metadata = item.model_dump(mode="json", exclude_none=True)
    source_note = render_frontmatter(metadata, f"# {item.title}\n\n{text}")
    candidate = CandidatePatch(
        id=source_id,
        type="candidate-patch",
        title=f"Ingest source: {item.title}",
        status="pending-review",
        sensitivity=sensitivity,
        created_at=instant,
        updated_at=instant,
        target_path=source_item_relative.as_posix(),
        content=source_note,
        expected_base_hash=None,
    )
    candidate_relative = Path(".constellation/candidates") / f"{candidate.id}.json"
    atomic_write_text(vault, candidate_relative, candidate.model_dump_json(indent=2) + "\n")
    manifest = {
        "schema_version": "0.1",
        "mode": "deferred-canonical",
        "source_id": source_id,
        "source_hash": digest,
        "source_size_bytes": len(data),
        "media_type": media_type,
        "ingested_at": instant.isoformat(),
        "input_path": relative.as_posix(),
        **({"source_url": source_url} if source_url else {}),
        **({"business_card": business_card} if business_card is not None else {}),
        "preserved_path": preserved_relative.as_posix(),
        "text_path": text_relative.as_posix(),
        "source_item_path": source_item_relative.as_posix(),
        "candidate_id": candidate.id,
        "candidate_path": candidate_relative.as_posix(),
        "registration": {
            "mode": registration_mode,
            "status": "pending-review" if registration_mode == "review" else "pending-automatic",
        },
        "extraction": extraction,
    }
    atomic_write_text(vault, manifest_relative, json.dumps(manifest, indent=2, sort_keys=True) + "\n")
    promotion: dict[str, str] | None = None
    if registration_mode == "automatic":
        from .review import promote_candidate

        promotion = promote_candidate(
            vault,
            candidate.id,
            confirm=True,
            expected_base_hash=None,
        )
        manifest["registration"] = {"mode": "automatic", "status": "canonical"}
        atomic_write_text(
            vault,
            manifest_relative,
            json.dumps(manifest, indent=2, sort_keys=True) + "\n",
        )
    result = {
        "status": "registered" if promotion else "staged",
        **{
            key: str(value)
            for key, value in manifest.items()
            if key not in {"extraction", "registration", "business_card"}
        },
        **(
            {"business_card_fields": str(len(business_card["fields"]))}
            if business_card is not None
            else {}
        ),
    }
    if promotion:
        result["index_generation"] = promotion["index_generation"]
    result["manifest_path"] = manifest_relative.as_posix()
    result["extraction_status"] = str(extraction["status"])
    return result
